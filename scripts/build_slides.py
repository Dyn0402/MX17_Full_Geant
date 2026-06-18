#!/usr/bin/env python3
"""build_slides.py — assemble the MX17 results deck (16:9 .pptx).

Embeds the slide figures (scripts/make_slides_figures.py) plus selected
existing PNGs from docs/. Upload the resulting .pptx to Google Drive with
conversion to Google Slides for the collaborator.

Deck layout
  Section A (INFN request): title, setup, production/reach, July optimistic,
    post-LS3.
  Section B (broader results, for D.N.): production budget, angular-resolution
    floor, the "money plot", resolution vs pair kinematics, the E0 channel,
    backgrounds/timing, summary.

Caveat boxes are added as SEPARATE text shapes (clearly labelled) so the
collaborator can delete them before the INFN talk.

Usage:  /home/dylan/PycharmProjects/nTof_x17/venv/bin/python scripts/build_slides.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

REPO = Path(__file__).resolve().parent.parent
SF = REPO / "docs/slides/figs"           # generated figures
AR = REPO / "docs/angular_resolution/figs"
RP = REPO / "docs/report/figs"
NH = REPO / "docs/neutron_histories/figs"
GEO = REPO / "scripts"
OUT = REPO / "docs/slides/MX17_results.pptx"

# 16:9
SW, SH = Inches(13.333), Inches(7.5)
NAVY = RGBColor(0x1F, 0x3A, 0x5F)
GREY = RGBColor(0x55, 0x55, 0x55)
RED = RGBColor(0xB0, 0x20, 0x20)
CAVEAT_BG = RGBColor(0xFF, 0xF4, 0xD6)
CAVEAT_BORDER = RGBColor(0xD0, 0xA0, 0x30)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def _txbox(slide, l, t, w, h, lines, *, size=18, bold=False, color=GREY,
           align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.05):
    """lines: str or list of (text, dict-overrides) or plain str."""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        ov = {}
        if isinstance(ln, tuple):
            ln, ov = ln
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ov.get("align", align)
        p.line_spacing = line_spacing
        if ov.get("space_before"):
            p.space_before = Pt(ov["space_before"])
        r = p.add_run()
        r.text = ln
        f = r.font
        f.size = Pt(ov.get("size", size))
        f.bold = ov.get("bold", bold)
        f.color.rgb = ov.get("color", color)
        if ov.get("italic"):
            f.italic = True
    return box


def _title(slide, text, sub=None):
    bar = slide.shapes.add_shape(1, 0, 0, SW, Inches(0.12))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    _txbox(slide, Inches(0.5), Inches(0.22), Inches(12.9), Inches(0.7),
           text, size=26, bold=True, color=NAVY)
    if sub:
        _txbox(slide, Inches(0.5), Inches(0.95), Inches(12.9), Inches(0.5),
               sub, size=15, color=GREY)


def _img_fit(slide, path, l, t, max_w, max_h):
    """Add image scaled to fit (max_w,max_h) box, centred in it."""
    with Image.open(path) as im:
        iw, ih = im.size
    ar = iw / ih
    box_ar = max_w / max_h
    if ar > box_ar:
        w = max_w; h = Emu(int(max_w / ar))
    else:
        h = max_h; w = Emu(int(max_h * ar))
    l2 = Emu(int(l + (max_w - w) / 2))
    t2 = Emu(int(t + (max_h - h) / 2))
    return slide.shapes.add_picture(str(path), l2, t2, w, h)


def _caveat(slide, l, t, w, h, lines, header="Caveats / assumptions "
            "(separate box — delete before INFN if desired)"):
    box = slide.shapes.add_shape(1, l, t, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = CAVEAT_BG
    box.line.color.rgb = CAVEAT_BORDER; box.line.width = Pt(1)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(8); tf.margin_right = Pt(8); tf.margin_top = Pt(5)
    p0 = tf.paragraphs[0]; r = p0.add_run(); r.text = header
    r.font.size = Pt(11); r.font.bold = True
    r.font.color.rgb = RGBColor(0x8A, 0x6D, 0x10)
    for ln in lines:
        p = tf.add_paragraph(); p.line_spacing = 1.0
        rr = p.add_run(); rr.text = "•  " + ln
        rr.font.size = Pt(11); rr.font.color.rgb = RGBColor(0x6B, 0x55, 0x10)
    return box


def _redbox(slide, l, t, w, h, lines, header):
    """Prominent red caveat box (must be KEPT — unlike the gold delete-me box)."""
    box = slide.shapes.add_shape(1, l, t, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xFD, 0xE8, 0xE8)
    box.line.color.rgb = RGBColor(0xC0, 0x30, 0x30); box.line.width = Pt(1.75)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(9); tf.margin_right = Pt(9); tf.margin_top = Pt(6)
    p0 = tf.paragraphs[0]; r = p0.add_run(); r.text = header
    r.font.size = Pt(12.5); r.font.bold = True
    r.font.color.rgb = RGBColor(0xB0, 0x20, 0x20)
    for ln in lines:
        p = tf.add_paragraph(); p.line_spacing = 1.0; p.space_before = Pt(3)
        rr = p.add_run(); rr.text = "•  " + ln
        rr.font.size = Pt(11.5); rr.font.color.rgb = RGBColor(0x5A, 0x10, 0x10)
    return box


def _table(slide, l, t, w, h, rows):
    """rows: list of row-lists of cell strings; row 0 is the header."""
    nr, nc = len(rows), len(rows[0])
    gf = slide.shapes.add_table(nr, nc, l, t, w, h)
    tbl = gf.table
    for ci, frac in enumerate([0.28, 0.18, 0.21, 0.16, 0.17][:nc]):
        tbl.columns[ci].width = Emu(int(w * frac))
    for ri, row in enumerate(rows):
        for ci, txt in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.margin_left = Pt(6); cell.margin_right = Pt(6)
            cell.margin_top = Pt(2); cell.margin_bottom = Pt(2)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if ci else PP_ALIGN.LEFT
            r = p.add_run(); r.text = txt
            r.font.size = Pt(13 if ri else 12.5)
            r.font.bold = bool(ri == 0)
            r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if ri == 0 else NAVY
            cell.fill.solid()
            cell.fill.fore_color.rgb = (
                NAVY if ri == 0 else
                RGBColor(0xEE, 0xF2, 0xF7) if ri % 2 else
                RGBColor(0xFF, 0xFF, 0xFF))
    return gf


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def add(title=None, sub=None):
    s = prs.slides.add_slide(BLANK)
    if title:
        _title(s, title, sub)
    return s


# ════════════════════════════════════════════════════════════════════════════
# SECTION A — the INFN request
# ════════════════════════════════════════════════════════════════════════════

# 1 — TITLE -------------------------------------------------------------------
s = add()
band = s.shapes.add_shape(1, 0, Inches(2.2), SW, Inches(2.9))
band.fill.solid(); band.fill.fore_color.rgb = NAVY; band.line.fill.background()
_txbox(s, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.4),
       "The X17 search at n_TOF EAR2", size=42, bold=True,
       color=RGBColor(0xFF, 0xFF, 0xFF))
_txbox(s, Inches(0.8), Inches(3.7), Inches(11.7), Inches(1.0),
       "Full Geant4 simulation: production, acceptance, and the path to a "
       "measurement\nFinal July setup  ·  expected results  ·  the post-LS3 "
       "upgrade", size=19, color=RGBColor(0xD7, 0xE3, 0xF4))
_txbox(s, Inches(0.8), Inches(5.4), Inches(11.7), Inches(0.6),
       "D. Neff — 16 June 2026   |   simulation + analysis: this repository "
       "(MX17_Full_Geant)", size=14, color=GREY)
_notes(s, "Deck for INFN presentation of the July run plan + post-LS3 "
       "expectations, plus a broader set of simulation results. Section A "
       "(slides 2-5) is the INFN-requested core; Section B is supporting "
       "detail. All numbers from the full Geant4 campaign (5e8 neutrons, "
       "1e7 pair events).")

# 2 — SETUP -------------------------------------------------------------------
s = add("The experiment: a 4-arm e⁺ e⁻ pair spectrometer",
        "³He(n,γ)⁴He* → (X17 or IPC) → e⁺ e⁻ ; "
        "measure the opening angle in 4 Micromegas + scintillator arms")
_img_fit(s, GEO / "mx17_geometry_topdown.png",
         Inches(0.3), Inches(1.55), Inches(7.4), Inches(5.7))
_txbox(s, Inches(7.9), Inches(1.7), Inches(5.1), Inches(5.3), [
    ("Target", {"size": 18, "bold": True, "color": NAVY}),
    "³He gas capsule, 500 bar, r = 10 mm, 60 mm on axis; thin Al (0.6 mm) "
    "+ CFRP (0.9 mm) wall.",
    ("Tracking", {"size": 18, "bold": True, "color": NAVY,
                  "space_before": 8}),
    "4 Micromegas arms, front face ≈ 25 cm from the target, 30 mm drift "
    "gap — the e⁺ e⁻ directions.",
    ("Calorimetry / trigger", {"size": 18, "bold": True, "color": NAVY,
                               "space_before": 8}),
    "Liquid-scintillator + plastic layers per arm (assumed saturated/unusable "
    "for the high-energy measurement — tracking-only baseline).",
    ("The observable", {"size": 18, "bold": True, "color": NAVY,
                        "space_before": 8}),
    "X17 (m ≈ 17 MeV) makes a sharp opening-angle shoulder at θ ≈ "
    "109°; the IPC continuum is the irreducible background.",
], size=15, color=GREY)
_notes(s, "Geometry from scripts/plot_geometry.py (SimConfig defaults). Beam "
       "into the page. Production happens in the He-3 gas; the two leptons exit "
       "through the capsule wall (the dominant resolution limit, see later) and "
       "are tracked in the Micromegas. Calorimetry is conservatively treated as "
       "unusable at high energy, so the measurement rests on MM tracking only.")

# 3 — PRODUCTION / REACH ------------------------------------------------------
s = add("Production is settled — and it lives at MeV energies",
        "Directly counted ³He(n,γ) from 5×10⁸ neutrons; "
        "the γ-flash sets the high-energy reach")
_img_fit(s, SF / "fig_production_windows.png",
         Inches(0.3), Inches(1.55), Inches(8.2), Inches(5.7))
_txbox(s, Inches(8.7), Inches(1.7), Inches(4.4), Inches(5.4), [
    ("Sub-keV is dead.", {"size": 16, "bold": True, "color": RED}),
    "The gas is opaque to (n,p); radiative capture is capped at "
    "σₙγ/σₙₚ ≈ 10⁻⁸ → ~0.1 "
    "X17/day.",
    ("The MeV window is the measurement.", {"size": 16, "bold": True,
                                            "color": NAVY, "space_before": 10}),
    "Above ~100 keV the gas is thin and the rate table is confirmed by Geant4 "
    "(within 10–70%). ~33 X17/day produced over the full range.",
    ("The γ-flash limits the reach.", {"size": 16, "bold": True,
                                            "color": NAVY, "space_before": 10}),
    "Faster neutrons arrive sooner, closer to the flash. How early the "
    "detectors recover sets the max neutron energy — this is the LS3 "
    "lever.",
], size=14, color=GREY)
_notes(s, "From docs/report/mev_note.tex + analysis/mev. The per-bin X17/day "
       "uses alpha_IPC = 3.5e-3. Green = July reach (0.2-0.7 MeV), blue = the "
       "extra reach LS3 buys (0.7-2 MeV). Top axis is neutron TOF to EAR2 at "
       "19.5 m: 2 MeV = 1.0 us, 0.7 MeV = 1.7 us, 0.2 MeV = 3.2 us. Recorded "
       "X17 in 30 d: 64 (July) -> 220 (LS3).")

# 4 — JULY OPTIMISTIC ---------------------------------------------------------
s = add("July run: expected result (optimistic scenario)",
        "Window 0.2–0.7 MeV — conservative γ-flash recovery "
        "(detectors usable by ≈ 1.7 µs); 30-day run")
_img_fit(s, SF / "fig_stacked_july.png",
         Inches(0.25), Inches(1.5), Inches(8.0), Inches(4.6))
_txbox(s, Inches(8.4), Inches(1.6), Inches(4.7), Inches(3.3), [
    ("≈ 64 X17 on ≈ 3 100 IPC", {"size": 20, "bold": True,
                                           "color": NAVY}),
    "recorded in 30 days (S/B ≈ 0.021).",
    ("The peak is diluted, not visible by eye.", {"size": 15, "bold": True,
                                                  "color": RED,
                                                  "space_before": 8}),
    "Capsule multiple scattering smears θ by σ68 ≈ 14°; the "
    "109° shoulder (green dotted = truth) becomes a broad bump. The "
    "measurement is a template fit, not a cut-and-count.",
], size=14, color=GREY)
_caveat(s, Inches(0.4), Inches(6.25), Inches(12.5), Inches(1.05), [
    "Optimistic: best-estimator (target-centre chord) smearing; no pile-up, "
    "no γ-flash losses, no θ-dependent MM acceptance shaping.",
    "α_IPC = 3.5×10⁻³ (¹²C/⁸Be-anchored) and "
    "X17/IPC = 2.5% are assumed inputs (provenance pending). Table α = "
    "2.1×10⁻³ would scale signal ×0.6 (→ ~38 X17).",
    "At-rest kinematics (shoulder shifts a few degrees with the Eₙ boost). "
    "MM-double acceptance 19.6% (X17) / 23.6% (IPC).",
])
_notes(s, "Conservative July choice: gamma-flash recovery only by ~1.7 us, so "
       "max neutron energy ~0.7 MeV (this 'energy cutoff' framing is how the "
       "collaboration thinks). 143 captures in the window -> 64 recorded X17 / "
       "3093 IPC over 30 days at alpha=3.5e-3. The caveat box is a separate "
       "shape: delete it for the INFN talk if desired. Figure: "
       "scripts/make_slides_figures.py.")

# 5 — POST-LS3 ----------------------------------------------------------------
s = add("After LS3: extend the reach to 2 MeV neutrons",
        "Faster γ-flash recovery (≈ 1.0 µs) → window 0.2–2 MeV; "
        "same analysis, ×3.4 the statistics")
_img_fit(s, SF / "fig_stacked_compare.png",
         Inches(0.2), Inches(1.6), Inches(12.9), Inches(4.55))
_txbox(s, Inches(0.5), Inches(6.2), Inches(12.4), Inches(1.1), [
    ("July 64 X17 → post-LS3 220 X17 recorded (30 d), on the same "
     "S/B ≈ 0.021.", {"size": 17, "bold": True, "color": NAVY}),
    "Reaching 2 MeV neutrons (faster γ-flash recovery) adds the most "
    "productive energy decade. The spectrum shape is unchanged — the gain "
    "is purely statistical, which is exactly what a template fit converts into "
    "significance.",
], size=14, color=GREY)
_notes(s, "LS3 goal stated by the collaboration: achieve neutron energy of at "
       "least 2 MeV. Same plot as July, window extended 0.7->2 MeV. Recorded "
       "X17 30 d: 220 vs 64 (x3.4). S/B identical (set by X17/IPC branching x "
       "acceptance ratio, window-independent). The decisive question remains "
       "the gamma-flash recovery time, measured on the real setup.")

# ════════════════════════════════════════════════════════════════════════════
# SECTION B — broader results
# ════════════════════════════════════════════════════════════════════════════

# 6 — SECTION DIVIDER ---------------------------------------------------------
s = add()
band = s.shapes.add_shape(1, 0, Inches(3.0), SW, Inches(1.5))
band.fill.solid(); band.fill.fore_color.rgb = NAVY; band.line.fill.background()
_txbox(s, Inches(0.8), Inches(3.25), Inches(11.7), Inches(1.0),
       "Supporting results", size=36, bold=True,
       color=RGBColor(0xFF, 0xFF, 0xFF))
_txbox(s, Inches(0.8), Inches(4.7), Inches(11.7), Inches(0.6),
       "Production budget · the angular-resolution floor · the analysis "
       "strategy · the γ-dark E0 channel · backgrounds", size=16,
       color=GREY)

# 6b — GEOMETRY: SIDE VIEW + 3D ----------------------------------------------
s = add("Detector geometry — side view and 3D",
        "Companion views to the top-down layout: the ³He capsule shape, "
        "beam-end-on incidence, and the full 4-arm arrangement")
_img_fit(s, GEO / "mx17_geometry_sideview.png",
         Inches(0.25), Inches(1.55), Inches(7.1), Inches(5.4))
_img_fit(s, GEO / "mx17_geometry_3d.png",
         Inches(7.5), Inches(1.55), Inches(5.6), Inches(5.4))
_txbox(s, Inches(0.4), Inches(7.0), Inches(6.7), Inches(0.4),
       "Side view (XY): beam enters end-on along +Y into the capsule; "
       "arms 0/1 (±X) shown.", size=12, color=GREY)
_txbox(s, Inches(7.5), Inches(7.0), Inches(5.6), Inches(0.4),
       "3D: 4 arms (±X, ±Z), each MM + scintillator stack around the target.",
       size=12, color=GREY)
_notes(s, "Backup geometry slide (scripts/plot_geometry.py). Side view shows "
       "the capsule bottle profile (20 mm dia x 40 mm + domed ends, neck/valve "
       "on axis) and the beam hitting it end-on along +Y. The 3D isometric "
       "shows all four arms; arms 2,3 are along +/-Z (not shown in the side "
       "view).")

# 7 — ANGULAR RESOLUTION FLOOR ------------------------------------------------
s = add("The angular resolution is hardware-frozen",
        "Opening-angle resolution is set by multiple scattering in the "
        "capsule wall — no reconstruction can beat it")
_img_fit(s, AR / "fig_ms_budget.png",
         Inches(0.3), Inches(1.6), Inches(7.6), Inches(5.5))
_txbox(s, Inches(8.1), Inches(1.7), Inches(4.9), Inches(5.4), [
    ("σ68(Δθ) ≈ 14.5°", {"size": 26, "bold": True,
                                                 "color": NAVY}),
    "full accepted spectrum (≈ 11.5–12.5° for symmetric pairs), "
    "even assuming a perfect vertex.",
    ("Why:", {"size": 16, "bold": True, "color": NAVY, "space_before": 10}),
    "The 0.6 mm Al + 0.9 mm CFRP wall (~11 mm from the vertex) is 79% of the "
    "material budget. A scatter that early rotates the whole visible track.",
    ("Verdict:", {"size": 16, "bold": True, "color": RED, "space_before": 10}),
    "All 7 reconstruction methods land within 2° of each other. The only "
    "hardware lever is an all-composite vessel (→ ~10°). The recovery "
    "path is statistical: fit the shape.",
], size=14, color=GREY)
_notes(s, "From docs/angular_resolution/angular_resolution_note. Highland "
       "budget on the new STEP capsule. The vertex-constraint trick "
       "(target-centre chord) is already within 0.5 deg of the perfect-vertex "
       "oracle because the small target banks the gain. Detector spatial "
       "resolution costs <0.5 deg. This is the single most important hardware "
       "fact for the measurement.")

# 8 — THE MONEY PLOT ----------------------------------------------------------
s = add("Why the analysis is a template fit, not a peak cut",
        "Truth → stacked on IPC → both smeared: the shoulder washes "
        "out, but its shape is known to ~1% from 10⁷ simulated pairs")
_img_fit(s, AR / "fig_theta_money.png",
         Inches(0.25), Inches(1.6), Inches(12.85), Inches(4.7))
_txbox(s, Inches(0.5), Inches(6.35), Inches(12.4), Inches(0.95), [
    "Cut-and-count on θ is dead (recorded S/B ≈ 0.02). What survives: "
    "a fit of the known smeared-signal template + the smooth IPC continuum, "
    "with resolution entering as dilution of statistical power rather than as "
    "a cliff. This is the natural go/no-go calculation.",
], size=14, color=GREY)
_notes(s, "The 'money plot' from the angular note. Per-leg smearing drawn from "
       "the measured P(psi|KE) response (best estimator). Normalised to "
       "full-production 30-day recorded yields (191 X17 / 9190 IPC at "
       "alpha=2.1e-3). The middle vs right panels show the shoulder before/after "
       "smearing on the IPC continuum.")

# 9 — RESOLUTION VS KINEMATICS ------------------------------------------------
s = add("Resolution is a soft-track problem",
        "σ68 vs pair symmetry, and the +bias (correctable in a "
        "response-matrix fit) — X17 is sharpest at its symmetric 109° shoulder")
_img_fit(s, AR / "fig_sigma68_vs_minke.png",
         Inches(0.4), Inches(1.7), Inches(12.5), Inches(5.0))
_txbox(s, Inches(0.5), Inches(6.75), Inches(12.4), Inches(0.55),
       "Below ~3 MeV a leg is directionless after the wall (σ68 → "
       "40°); above ~8 MeV it plateaus at ~12°. X17 is sharpest at "
       "its symmetric 109° shoulder — the good news.",
       size=14, color=GREY)
_notes(s, "X17 has KE- + KE+ ~ 19.6 MeV so one leg is always <= 9.8 MeV; the "
       "resolution is driven by the softer leg. MM-internal quality cuts "
       "(straightness, dE/dx) could reject soft legs and recover ~12.5 deg "
       "without a calorimeter -- the one open reconstruction study.")

# 10 — E0 CHANNEL -------------------------------------------------------------
s = add("The γ-dark E0 channel: completeness, not a game-changer",
        "0⁺→0⁺ captures make no photon but can make a massive "
        "X17 — yet the rate is tiny and sub-keV")
_img_fit(s, RP / "fig_e0_final_metric.png",
         Inches(0.3), Inches(1.6), Inches(8.3), Inches(5.4))
_txbox(s, Inches(8.7), Inches(1.7), Inches(4.4), Inches(5.4), [
    ("E0 is real but small.", {"size": 16, "bold": True, "color": NAVY}),
    "The 0⁺ doorway sits sub-keV, where the same 10⁻⁸ radiative "
    "branching that kills the thermal program applies — times f ≤ "
    "10⁻².",
    ("The flash readout cannot see it.", {"size": 16, "bold": True,
                                          "color": RED, "space_before": 10}),
    "Sub-keV → TOF > 45 µs → arrives after the window closes. "
    "Only a dedicated thermal trigger records it (≤ 0.06 X17/day).",
    ("So:", {"size": 16, "bold": True, "color": NAVY, "space_before": 10}),
    "Leaving E0 out does not bias the baseline high-energy analysis. Worth "
    "adding to the generator for completeness only.",
], size=14, color=GREY)
_notes(s, "From docs/report/e0_final_metric_note + e0_branch docs. Addresses "
       "Alberto's E0 question. Key physics: 0+->0+ is photon-dark because the "
       "photon is massless, but a massive vector/scalar X17 is allowed. The "
       "rate is bounded by the (tiny) E0 capture rate. Left panel: IPC "
       "background pairs/day; right: X17 signal. Gold/blue = thermal vs flash "
       "trigger regions.")

# 10b — EXPECTED SIGNIFICANCE (statistical ceiling) --------------------------
s = add("Expected significance — a statistical ceiling, not the final CL",
        "Profile-likelihood (template-fit) Asimov significance, with the IPC "
        "background shape & normalisation taken as exactly known")
_table(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.45), [
    ["Scenario (30-day run)", "Recorded  S / B",
     "Shape-fit Z  (smeared, realistic)", "Z if no smearing",
     "Naïve S/√B  (whole spectrum)"],
    ["July  ·  0.2–0.7 MeV", "64 / 3 093", "2.6 σ", "3.9 σ", "1.2"],
    ["Post-LS3  ·  0.2–2 MeV", "220 / 10 600", "4.9 σ", "7.3 σ", "2.1"],
])
_txbox(s, Inches(0.5), Inches(3.05), Inches(5.85), Inches(4.1), [
    ("What IS in these numbers", {"size": 15, "bold": True, "color": NAVY}),
    "Binned profile-likelihood ratio (Cowan et al. Asimov formula) — the "
    "rigorous form of “count in a region”, optimally weighted over "
    "all bins.",
    "Same templates as slides 4–5: best-estimator (target-centre chord) "
    "capsule multiple-scattering smearing; α_IPC = 3.5×10⁻³, X17/IPC = 2.5%, "
    "MM-double acceptance.",
    "Binning-independent: 2° and 8° (×4) bins agree to <1%.",
    "Background normalisation can be floated from the θ ≲ 90° region (X17-free) "
    "to ~2% — negligible. The shape fit ~doubles the naïve whole-spectrum "
    "S/√B by weighting the high-θ shoulder.",
], size=12, color=GREY, line_spacing=1.03)
_redbox(s, Inches(6.55), Inches(3.05), Inches(6.28), Inches(4.1), [
    "The IPC angular SHAPE is assumed equal to the simulation truth; its "
    "uncertainty — which will dominate the real CL — is NOT propagated.",
    "Physics: a single E1 transition is assumed (probably fair, unquantified "
    "here).",
    "Detector/acceptance shape distortion is omitted: poor separation of "
    "near-collinear tracks (small θ), and finite active area shaping the catch "
    "probability for wide-angle pairs (large θ) — the signal region.",
    "With B ≈ 3 000–10 000, a few-% shape error in the high-θ signal region "
    "already rivals the statistical fluctuation — by post-LS3 the IPC shape, "
    "not statistics, sets the achievable CL.",
    "Read 2.6σ / 4.9σ as a best-case ceiling. A defensible CL needs the "
    "MEASURED IPC shape + full detector response folded in (open work).",
], header="⚠  Why this is a ceiling, not a confidence level  (keep this box)")
_notes(s, "Significance projection (open item from the summary slide). Numbers "
       "from scripts/make_slides_figures.py templates fed through an Asimov "
       "binned profile-likelihood: Z = sqrt(2 sum[(s+b)ln(1+s/b)-s]). July "
       "smeared Z=2.64 (truth 3.94), LS3 smeared Z=4.89 (truth 7.29); 2deg vs "
       "8deg binning agree to <1%; a best single theta-window cut-and-count "
       "gives ~2.5 / 4.7 sigma (the broad shoulder makes counting nearly as "
       "good as the full fit). CRITICAL CAVEAT (the red box): these assume the "
       "IPC angular shape equals MC truth. Not included: physics shape "
       "uncertainty (single-E1 assumption), and — emphasised by D.N. — detector "
       "/ acceptance shape distortion: track-pair separation inefficiency at "
       "small opening angle, and active-area geometric acceptance falling off "
       "for wide-angle pairs (the signal region). These reshape the IPC "
       "template where the signal sits and will dominate the final CL. Do not "
       "quote these as a final confidence level.")

# 11 — SUMMARY ----------------------------------------------------------------
s = add("Summary & status")
_txbox(s, Inches(0.7), Inches(1.6), Inches(12.0), Inches(5.4), [
    ("Production is settled.", {"size": 20, "bold": True, "color": NAVY}),
    "Full Geant4 (5×10⁸ neutrons) confirms the rate table where it is "
    "valid: ~33 X17/day produced, living at MeV energies. The sub-keV "
    "measurement is dead (self-shielding).",
    ("July (current hardware): ~64 recorded X17 / 30 d.", {"size": 20,
        "bold": True, "color": NAVY, "space_before": 12}),
    "Window 0.2–0.7 MeV under a conservative γ-flash recovery; on "
    "~3 100 IPC (S/B ≈ 0.02). Extractable by a template fit, not a peak "
    "cut.",
    ("Post-LS3: ~220 recorded X17 / 30 d.", {"size": 20, "bold": True,
        "color": NAVY, "space_before": 12}),
    "Reaching 2 MeV neutrons (×3.4 statistics). The decisive measurement "
    "is the real γ-flash recovery time.",
    ("The resolution is hardware-frozen at σ68 ≈ 14°.",
     {"size": 20, "bold": True, "color": NAVY, "space_before": 12}),
    "Capsule wall MS; an all-composite vessel is the only lever. E0 is a "
    "completeness item, not a driver.",
    ("Open items:", {"size": 16, "bold": True, "color": RED,
                     "space_before": 12}),
    "α_IPC & X17/IPC provenance (Alberto); Eₙ-dependent generator "
    "kinematics; the template-fit significance projection; γ-flash "
    "recovery measurement.",
], size=15, color=GREY)
_notes(s, "Wrap-up. The program is now an acceptance-and-backgrounds problem, "
       "not a production one. Next analysis step: wire the smeared templates "
       "into the fast-MC likelihood fit and quote expected significance vs run "
       "time.")

# ════════════════════════════════════════════════════════════════════════════
# BACKUP — factor-4 rebinned spectra (colleague request)
# ════════════════════════════════════════════════════════════════════════════

# B1 — BACKUP DIVIDER ---------------------------------------------------------
s = add()
band = s.shapes.add_shape(1, 0, Inches(3.0), SW, Inches(1.5))
band.fill.solid(); band.fill.fore_color.rgb = NAVY; band.line.fill.background()
_txbox(s, Inches(0.8), Inches(3.25), Inches(11.7), Inches(1.0),
       "Backup: factor-4 rebinned spectra", size=34, bold=True,
       color=RGBColor(0xFF, 0xFF, 0xFF))
_txbox(s, Inches(0.8), Inches(4.7), Inches(11.7), Inches(0.7),
       "The slides 4–5 stacked spectra with 4 base bins merged into one "
       "(2° → 8° opening-angle bins) — same events, coarser binning.",
       size=16, color=GREY)
_notes(s, "Backup versions of the slide 4 (July) and slide 5 (July vs LS3) "
       "stacked opening-angle spectra, rebinned by a factor of 4 (2 deg -> 8 "
       "deg bins) per collaborator request. Identical event samples, yields and "
       "normalisation as the main slides; only the histogram binning differs. "
       "90 base bins are not divisible by 4, so the final 176-180 deg bin holds "
       "the 2 leftover base bins (negligible -- the distribution is empty "
       "there). Figures: scripts/make_slides_figures.py (rebin=4).")

# B2 — JULY, REBINNED ---------------------------------------------------------
s = add("July run, rebinned ×4 (8° bins)",
        "Same as slide 4 — window 0.2–0.7 MeV, 30-day run — with 4 base "
        "bins merged into one")
_img_fit(s, SF / "fig_stacked_july_rebin4.png",
         Inches(0.25), Inches(1.5), Inches(8.0), Inches(4.6))
_txbox(s, Inches(8.4), Inches(1.6), Inches(4.7), Inches(4.6), [
    ("≈ 64 X17 on ≈ 3 100 IPC", {"size": 20, "bold": True, "color": NAVY}),
    "recorded in 30 days (S/B ≈ 0.021) — yields unchanged from slide 4.",
    ("Coarser binning, same statistics.", {"size": 15, "bold": True,
                                            "color": RED, "space_before": 8}),
    "Merging 4 base bins (2° → 8°) smooths the per-bin fluctuations and "
    "makes the smeared 109° excess (red, over green-dotted truth) easier "
    "to read by eye — but it adds no information. Significance is set by the "
    "template fit, which is binning-independent.",
], size=14, color=GREY)
_notes(s, "Backup rebinned (x4) twin of slide 4. Same recorded yields (64 X17 / "
       "3093 IPC, S/B 0.021); only the display binning is coarser. Use to "
       "answer the collaborator's rebinning request; do not read a confidence "
       "level off the binning -- that comes from the template-fit / "
       "counting analysis (see notes).")

# B3 — JULY vs LS3, REBINNED --------------------------------------------------
s = add("July vs post-LS3, rebinned ×4 (8° bins)",
        "Same as slide 5 — the two windows side by side — with 4 base "
        "bins merged into one")
_img_fit(s, SF / "fig_stacked_compare_rebin4.png",
         Inches(0.2), Inches(1.6), Inches(12.9), Inches(4.55))
_txbox(s, Inches(0.5), Inches(6.2), Inches(12.4), Inches(1.1), [
    ("July 64 X17 → post-LS3 220 X17 recorded (30 d), same "
     "S/B ≈ 0.021.", {"size": 17, "bold": True, "color": NAVY}),
    "Rebinned ×4 (2° → 8°) for readability; the shape and yields are "
    "identical to slide 5. The post-LS3 gain is purely statistical.",
], size=14, color=GREY)
_notes(s, "Backup rebinned (x4) twin of slide 5 (fig_stacked_compare_rebin4). "
       "Identical to slide 5 apart from the 8 deg display binning.")

prs.save(str(OUT))
print(f"wrote {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
